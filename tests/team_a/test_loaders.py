import pytest

from vector_db import loaders, ingestion


class FakeCollection:
    def __init__(self):
        self.added = None

    def upsert(self, ids, documents, metadatas):
        self.added = (ids, documents, metadatas)


def test_unsupported_extension_raises(tmp_path):
    f = tmp_path / "data.csv"
    f.write_text("a,b,c")
    with pytest.raises(ValueError):
        loaders.load_text(str(f))


def test_missing_file_raises():
    with pytest.raises(ValueError):
        loaders.load_text("nope.md")


def test_old_ppt_format_raises(tmp_path):
    f = tmp_path / "deck.ppt"
    f.write_bytes(b"\x00\x01")
    with pytest.raises(ValueError):
        loaders.load_text(str(f))


def test_markdown_loads_as_one_unit(tmp_path):
    f = tmp_path / "notes.md"
    f.write_text("# Title\n\nStudy content about mitochondria and ATP.", encoding="utf-8")

    units = loaders.load_text(str(f))

    assert len(units) == 1
    assert "mitochondria" in units[0].text
    assert units[0].source_file == "notes.md"
    assert units[0].slide is None
    assert units[0].transcribed is False


def test_empty_text_file_yields_no_units(tmp_path):
    f = tmp_path / "blank.txt"
    f.write_text("   \n  ", encoding="utf-8")
    assert loaders.load_text(str(f)) == []


def test_pptx_loads_slides_with_provenance(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Photosynthesis converts light to energy"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    units = loaders.load_text(str(path))

    assert len(units) >= 1
    assert units[0].slide == 1
    assert "Photosynthesis" in units[0].text


def test_ingest_file_stores_markdown_with_metadata(tmp_path):
    f = tmp_path / "notes.md"
    f.write_text("# Cells\n\nThe mitochondria is the powerhouse of the cell.", encoding="utf-8")
    coll = FakeCollection()

    n = ingestion.ingest_file(str(f), topic="biology", session_id="s1", collection=coll)

    assert n >= 1
    ids, docs, metas = coll.added
    assert len(ids) == n
    assert all(m["topic"] == "biology" for m in metas)
    assert all(m["session_id"] == "s1" for m in metas)
    assert all(m["source_file"] == "notes.md" for m in metas)
    assert all(m["transcribed"] is False for m in metas)


def test_ingest_file_routes_text_pdf_to_existing_path(tmp_path, monkeypatch):
    pdf = tmp_path / "notes.pdf"
    pdf.write_bytes(b"%PDF-1.4 dummy")
    called = {}

    def fake_ingest_pdf(path, topic, session_id, collection=None):
        called["path"] = path
        return 7

    monkeypatch.setattr(loaders, "pdf_has_text_layer", lambda p, **k: True)
    monkeypatch.setattr(ingestion, "ingest_pdf", fake_ingest_pdf)
    result = ingestion.ingest_file(str(pdf), topic="t", session_id="s1")

    assert result == 7
    assert called["path"] == str(pdf)


def test_pdf_has_text_layer_true_for_text_pdf(tmp_path):
    import fitz

    p = tmp_path / "text.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "This page carries a real text layer with plenty of readable words.")
    doc.save(str(p))
    doc.close()

    assert loaders.pdf_has_text_layer(str(p)) is True


def test_pdf_has_text_layer_false_for_blank(tmp_path):
    import fitz

    p = tmp_path / "blank.pdf"
    doc = fitz.open()
    doc.new_page()
    doc.save(str(p))
    doc.close()

    assert loaders.pdf_has_text_layer(str(p)) is False


def test_image_pdf_routes_to_ocr_and_flags_transcribed(tmp_path, monkeypatch):
    pdf = tmp_path / "scan.pdf"
    pdf.write_bytes(b"%PDF-1.4 dummy")
    monkeypatch.setattr(loaders, "pdf_has_text_layer", lambda p, **k: False)
    monkeypatch.setattr(
        loaders,
        "load_pdf_ocr",
        lambda path, source: [
            loaders.TextUnit(text="handwritten note about cells", source_file=source, page=1, transcribed=True)
        ],
    )
    coll = FakeCollection()

    n = ingestion.ingest_file(str(pdf), topic="bio", session_id="s1", collection=coll)

    assert n >= 1
    _, _, metas = coll.added
    assert all(m["transcribed"] is True for m in metas)
    assert all(m.get("page") == 1 for m in metas)


def test_ocr_image_dispatch_selects_provider(monkeypatch):
    monkeypatch.setattr(loaders.settings, "OCR_PROVIDER", "tesseract")
    monkeypatch.setattr(loaders, "_ocr_tesseract", lambda b: "TESS")
    assert loaders._ocr_image(b"x") == "TESS"

    monkeypatch.setattr(loaders.settings, "OCR_PROVIDER", "groq")
    monkeypatch.setattr(loaders, "_ocr_groq", lambda b: "GROQ")
    assert loaders._ocr_image(b"x") == "GROQ"


def test_invalid_ocr_provider_raises(monkeypatch):
    monkeypatch.setattr(loaders.settings, "OCR_PROVIDER", "nope")
    with pytest.raises(ValueError):
        loaders._ocr_image(b"x")


def test_load_pdf_ocr_sets_transcribed_with_page(monkeypatch):
    monkeypatch.setattr(loaders, "_render_pdf_pages", lambda path, dpi=200: [b"p1", b"p2"])
    monkeypatch.setattr(loaders, "_ocr_image", lambda png: "text " + png.decode())

    units = loaders.load_pdf_ocr("scan.pdf", "scan.pdf")

    assert len(units) == 2
    assert all(u.transcribed for u in units)
    assert [u.page for u in units] == [1, 2]


def test_segments_to_units_windows_with_timestamps():
    segments = [(0.0, "a" * 500), (5.0, "b" * 500), (12.0, "tail")]
    units = loaders._segments_to_units(segments, "lecture.mp3", window_chars=800)

    assert len(units) == 2
    assert units[0].timestamp == 0.0
    assert units[1].timestamp == 12.0
    assert all(u.transcribed for u in units)
    assert all(u.source_file == "lecture.mp3" for u in units)


def test_asr_dispatch_selects_provider(monkeypatch):
    monkeypatch.setattr(loaders.settings, "ASR_PROVIDER", "local")
    monkeypatch.setattr(loaders, "_asr_local", lambda p: [(0.0, "local text")])
    assert loaders._transcribe("a.mp3") == [(0.0, "local text")]

    monkeypatch.setattr(loaders.settings, "ASR_PROVIDER", "groq")
    monkeypatch.setattr(loaders, "_asr_groq", lambda p: [(1.0, "groq text")])
    assert loaders._transcribe("a.mp3") == [(1.0, "groq text")]


def test_invalid_asr_provider_raises(monkeypatch):
    monkeypatch.setattr(loaders.settings, "ASR_PROVIDER", "nope")
    with pytest.raises(ValueError):
        loaders._transcribe("a.mp3")


def test_audio_file_loads_as_transcribed_units(tmp_path, monkeypatch):
    audio = tmp_path / "lecture.mp3"
    audio.write_bytes(b"fake audio")
    monkeypatch.setattr(loaders, "_transcribe", lambda p: [(0.0, "intro to agents"), (8.0, "tool use")])

    units = loaders.load_text(str(audio))

    assert len(units) >= 1
    assert all(u.transcribed for u in units)
    assert units[0].timestamp == 0.0


def test_video_disabled_raises(tmp_path, monkeypatch):
    video = tmp_path / "lecture.mp4"
    video.write_bytes(b"fake video")
    monkeypatch.setattr(loaders.settings, "VIDEO_ENABLED", False)
    with pytest.raises(ValueError):
        loaders.load_text(str(video))


def test_video_enabled_routes_to_video_loader(tmp_path, monkeypatch):
    video = tmp_path / "lecture.mp4"
    video.write_bytes(b"fake video")
    monkeypatch.setattr(loaders.settings, "VIDEO_ENABLED", True)
    monkeypatch.setattr(
        loaders, "load_video",
        lambda path, source: [loaders.TextUnit(text="spoken content", source_file=source, timestamp=0.0, transcribed=True)],
    )

    units = loaders.load_text(str(video))

    assert len(units) == 1
    assert units[0].transcribed is True
