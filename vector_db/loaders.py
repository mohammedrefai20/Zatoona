import base64
import io
import os
from dataclasses import dataclass

from config import settings

TEXT_EXTS = {".md", ".txt"}
PPTX_EXTS = {".pptx"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a"}
VIDEO_EXTS = {".mp4"}
SUPPORTED = ".pdf .md .txt .pptx .mp3 .wav .m4a .mp4"

_OCR_PROMPT = (
    "Transcribe all text in this image exactly as written, preserving line breaks. "
    "Output only the transcribed text."
)


@dataclass
class TextUnit:
    text: str
    source_file: str
    page: int = None
    slide: int = None
    timestamp: float = None
    transcribed: bool = False


def load_text(path):
    if not os.path.isfile(path):
        raise ValueError(f"file not found: {path}")

    ext = os.path.splitext(path)[1].lower()
    source = os.path.basename(path)

    if ext in TEXT_EXTS:
        return _load_plain_text(path, source)
    if ext in PPTX_EXTS:
        return _load_pptx(path, source)
    if ext in AUDIO_EXTS:
        return load_audio(path, source)
    if ext in VIDEO_EXTS:
        if not settings.VIDEO_ENABLED:
            raise ValueError("video ingestion is disabled; set VIDEO_ENABLED=true to enable it")
        return load_video(path, source)
    if ext == ".ppt":
        raise ValueError("old .ppt format is not supported; save it as .pptx")
    raise ValueError(f"unsupported file type: {ext}. supported: {SUPPORTED}")


def _load_plain_text(path, source):
    with open(path, encoding="utf-8", errors="replace") as fh:
        text = fh.read().strip()
    if not text:
        return []
    return [TextUnit(text=text, source_file=source)]


def _load_pptx(path, source):
    from pptx import Presentation

    try:
        prs = Presentation(path)
    except Exception as exc:
        raise ValueError(f"could not read PowerPoint {source}: {exc}") from exc

    units = []
    for number, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append("Notes: " + notes)
        text = "\n".join(parts).strip()
        if text:
            units.append(TextUnit(text=text, source_file=source, slide=number))
    return units


def pdf_has_text_layer(path, min_chars_per_page=20):
    import fitz

    doc = fitz.open(path)
    pages = doc.page_count
    total = sum(len(page.get_text().strip()) for page in doc)
    doc.close()
    if pages == 0:
        return False
    return (total / pages) >= min_chars_per_page


def load_pdf_ocr(path, source):
    units = []
    for number, png in enumerate(_render_pdf_pages(path), start=1):
        text = (_ocr_image(png) or "").strip()
        if text:
            units.append(TextUnit(text=text, source_file=source, page=number, transcribed=True))
    return units


def _render_pdf_pages(path, dpi=200):
    import fitz

    doc = fitz.open(path)
    images = [page.get_pixmap(dpi=dpi).tobytes("png") for page in doc]
    doc.close()
    return images


def _ocr_image(png_bytes):
    provider = settings.OCR_PROVIDER
    if provider == "tesseract":
        return _ocr_tesseract(png_bytes)
    if provider == "groq":
        return _ocr_groq(png_bytes)
    if provider == "gemini":
        return _ocr_gemini(png_bytes)
    raise ValueError(f"invalid OCR_PROVIDER: {provider}")


def _ocr_tesseract(png_bytes):
    import pytesseract
    from PIL import Image

    return pytesseract.image_to_string(Image.open(io.BytesIO(png_bytes)))


def _ocr_groq(png_bytes):
    from groq import Groq

    if not settings.GROQ_API_KEY:
        raise RuntimeError("GROQ_API_KEY is not set")
    client = Groq(api_key=settings.GROQ_API_KEY)
    model = settings.VISION_MODEL or "meta-llama/llama-4-scout-17b-16e-instruct"
    b64 = base64.b64encode(png_bytes).decode()
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": _OCR_PROMPT},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                ],
            }
        ],
    )
    return resp.choices[0].message.content


def _ocr_gemini(png_bytes):
    from google import genai
    from google.genai import types

    if not settings.GEMINI_API_KEY:
        raise RuntimeError("GEMINI_API_KEY is not set")
    client = genai.Client(api_key=settings.GEMINI_API_KEY)
    model = settings.VISION_MODEL or "gemini-2.0-flash"
    resp = client.models.generate_content(
        model=model,
        contents=[_OCR_PROMPT, types.Part.from_bytes(data=png_bytes, mime_type="image/png")],
    )
    return resp.text


def load_audio(path, source):
    return _segments_to_units(_transcribe(path), source)


def load_video(path, source):
    audio_path = _extract_audio(path)
    try:
        return load_audio(audio_path, source)
    finally:
        if os.path.exists(audio_path):
            os.remove(audio_path)


def _segments_to_units(segments, source, window_chars=800):
    units = []
    buf = []
    buf_len = 0
    start_ts = None
    for start, text in segments:
        text = text.strip()
        if not text:
            continue
        if start_ts is None:
            start_ts = start
        buf.append(text)
        buf_len += len(text)
        if buf_len >= window_chars:
            units.append(TextUnit(text=" ".join(buf), source_file=source,
                                  timestamp=round(start_ts, 1), transcribed=True))
            buf, buf_len, start_ts = [], 0, None
    if buf:
        units.append(TextUnit(text=" ".join(buf), source_file=source,
                              timestamp=round(start_ts or 0.0, 1), transcribed=True))
    return units


def _transcribe(path):
    provider = settings.ASR_PROVIDER
    if provider == "local":
        return _asr_local(path)
    if provider == "groq":
        return _asr_groq(path)
    raise ValueError(f"invalid ASR_PROVIDER: {provider}")


def _asr_local(path):
    from faster_whisper import WhisperModel

    model = WhisperModel(settings.ASR_MODEL)
    segments, _ = model.transcribe(path)
    return [(seg.start, seg.text) for seg in segments]


def _asr_groq(path):
    from groq import Groq

    if not settings.GROQ_API_KEY:
        raise RuntimeError("GROQ_API_KEY is not set")
    client = Groq(api_key=settings.GROQ_API_KEY)
    with open(path, "rb") as fh:
        result = client.audio.transcriptions.create(
            file=(os.path.basename(path), fh.read()),
            model="whisper-large-v3-turbo",
            response_format="verbose_json",
        )
    segments = result.segments or []
    return [(s["start"], s["text"]) for s in segments]


def _extract_audio(video_path):
    import subprocess
    import tempfile

    import imageio_ffmpeg

    ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
    out = tempfile.mktemp(suffix=".wav")
    subprocess.run(
        [ffmpeg, "-y", "-i", video_path, "-vn", "-ac", "1", "-ar", "16000", out],
        check=True,
        capture_output=True,
    )
    return out
