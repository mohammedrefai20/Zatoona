import pytest

from vector_db import ingestion

MD = (
    "# Cells\n\n"
    "## Mitochondria\n\n"
    "The mitochondria is the powerhouse of the cell and produces ATP through respiration.\n"
)


class FakeCollection:
    def __init__(self):
        self.added = None

    def upsert(self, ids, documents, metadatas):
        self.added = (ids, documents, metadatas)


def test_ingest_markdown_stores_contextualized_with_metadata(tmp_path):
    f = tmp_path / "notes.md"
    f.write_text(MD, encoding="utf-8")
    coll = FakeCollection()

    n = ingestion.ingest_file(str(f), topic="biology", session_id="s1", collection=coll)

    assert n >= 1
    ids, docs, metas = coll.added
    assert len(ids) == n
    assert all(m["topic"] == "biology" for m in metas)
    assert all(m["session_id"] == "s1" for m in metas)
    assert all(m["source_file"] == "notes.md" for m in metas)
    assert all(m["transcribed"] is False for m in metas)
    assert any("headings" in m for m in metas)
    assert any("Mitochondria" in d for d in docs)
    assert all(i.startswith("s1:notes.md:") for i in ids)
    assert len(set(ids)) == len(ids)


def test_empty_document_stores_nothing(tmp_path):
    f = tmp_path / "blank.md"
    f.write_text("   \n", encoding="utf-8")
    coll = FakeCollection()
    assert ingestion.ingest_file(str(f), topic="t", session_id="s1", collection=coll) == 0


def test_unsupported_file_raises(tmp_path):
    f = tmp_path / "data.csv"
    f.write_text("a,b,c")
    with pytest.raises(ValueError):
        ingestion.ingest_file(str(f), topic="t", session_id="s1", collection=FakeCollection())


def test_missing_file_raises():
    with pytest.raises(ValueError):
        ingestion.ingest_file("does_not_exist.md", topic="t", session_id="s1")


def test_audio_routes_to_media_path_with_timestamps(tmp_path, monkeypatch):
    from vector_db import loaders

    audio = tmp_path / "lecture.mp3"
    audio.write_bytes(b"fake audio")
    monkeypatch.setattr(loaders, "_transcribe", lambda p: [(0.0, "intro to cells"), (8.0, "atp synthesis")])
    coll = FakeCollection()

    n = ingestion.ingest_file(str(audio), topic="bio", session_id="s1", collection=coll)

    assert n >= 1
    _, _, metas = coll.added
    assert all(m["transcribed"] is True for m in metas)
    assert all("timestamp" in m for m in metas)


def test_scanned_input_marks_transcribed(tmp_path, monkeypatch):
    from vector_db import chunking, docling_parser

    img = tmp_path / "scan.png"
    img.write_bytes(b"\x89PNG\r\n")
    monkeypatch.setattr(docling_parser, "parse", lambda p: "DLDOC")
    monkeypatch.setattr(docling_parser, "is_scanned", lambda p: True)
    monkeypatch.setattr(
        chunking, "chunk",
        lambda dl, source, transcribed=False: [
            chunking.Chunk(content="recognized text about cells", page=1, headings=[],
                           source_file=source, transcribed=transcribed)
        ],
    )
    coll = FakeCollection()

    n = ingestion.ingest_file(str(img), topic="bio", session_id="s1", collection=coll)

    assert n >= 1
    _, _, metas = coll.added
    assert all(m["transcribed"] is True for m in metas)
    assert all(m.get("page") == 1 for m in metas)


def test_ingest_pptx_stores_slide_provenance(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content (has a body placeholder)
    slide.shapes.title.text = "Photosynthesis"
    slide.placeholders[1].text = "Plants convert light energy into chemical energy in the chloroplasts."
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    coll = FakeCollection()
    try:
        n = ingestion.ingest_file(str(path), topic="bio", session_id="s1", collection=coll)
    except Exception as exc:  # Docling pptx backend unavailable
        pytest.skip(f"Docling pptx unavailable: {exc}")

    assert n >= 1
    _, _, metas = coll.added
    assert any("slide" in m for m in metas)


def test_scanned_pdf_ocr_end_to_end():
    pytest.skip("requires OCR models + image-PDF fixture; covered by quickstart.md Scenario 3")


def test_unknown_embedding_provider_raises(monkeypatch):
    from vector_db import embedder

    embedder._cached_ef = None
    monkeypatch.setattr(embedder.settings, "EMBEDDING_PROVIDER", "groq")
    with pytest.raises(RuntimeError):
        embedder.get_embedding_function()


def test_openai_provider_without_key_raises(monkeypatch):
    from vector_db import embedder

    embedder._cached_ef = None
    monkeypatch.setattr(embedder.settings, "EMBEDDING_PROVIDER", "openai")
    monkeypatch.setattr(embedder.settings, "OPENAI_API_KEY", "")
    with pytest.raises(RuntimeError):
        embedder.get_embedding_function()


def test_auto_provider_without_key_uses_local(monkeypatch):
    from vector_db import embedder

    embedder._cached_ef = None
    monkeypatch.setattr(embedder.settings, "EMBEDDING_PROVIDER", "auto")
    monkeypatch.setattr(embedder.settings, "OPENAI_API_KEY", "")
    monkeypatch.setattr(embedder, "_local_ef", lambda: "LOCAL_EF")
    assert embedder.get_embedding_function() == "LOCAL_EF"
    embedder._cached_ef = None
